#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import re
import json
import shutil
import threading
from datetime import datetime
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext


# ---------------------------------------------------------------------------
# 置換ロジック（1パス同時置換・冪等）
# ---------------------------------------------------------------------------

def _make_pat(old: str, fullhalf_sensitive: bool) -> str:
    if fullhalf_sensitive:
        return re.escape(old)
    parts = []
    for ch in old:
        code = ord(ch)
        if 0xFF01 <= code <= 0xFF5E:
            half = chr(code - 0xFEE0)
            parts.append(f'(?:{re.escape(ch)}|{re.escape(half)})')
        elif 0x21 <= code <= 0x7E:
            full = chr(code + 0xFEE0)
            parts.append(f'(?:{re.escape(ch)}|{re.escape(full)})')
        elif ch == ' ':
            parts.append(r'(?:　| )')
        elif ch == '　':
            parts.append(r'(?:　| )')
        else:
            parts.append(re.escape(ch))
    return ''.join(parts)


def build_combined_replacer(pairs_raw: list[tuple[str, str]],
                             case_sensitive: bool,
                             fullhalf_sensitive: bool):
    """
    全ペアを 1 本の正規表現で同時処理する冪等な置換関数を返す。
    置換後の値をガードパターンとして最優先登録することで、
    複数回実行しても結果が変わらない冪等性を保証する。
    """
    entries = []
    for old, new in pairs_raw:
        if old:
            entries.append((old, _make_pat(old, fullhalf_sensitive), new))

    if not entries:
        return None

    entries.sort(key=lambda e: len(e[0]), reverse=True)

    flags = re.UNICODE | (0 if case_sensitive else re.IGNORECASE)

    guard_values = sorted(
        set(new for _, new in pairs_raw if new),
        key=len, reverse=True,
    )
    guard_pats = [re.escape(v) for v in guard_values]
    guard_set  = {v.lower() for v in guard_values} if not case_sensitive else set(guard_values)

    individual = [(re.compile(f'(?:{pat})', flags), new) for _, pat, new in entries]
    all_pats   = guard_pats + [pat for _, pat, _ in entries]
    combined   = re.compile('|'.join(f'(?:{p})' for p in all_pats), flags)

    def replacer(text: str) -> str:
        def sub(m: re.Match) -> str:
            s     = m.group(0)
            s_key = s.lower() if not case_sensitive else s
            if s_key in guard_set:
                return s
            for irx, new in individual:
                if irx.fullmatch(s):
                    return new
            return s
        return combined.sub(sub, text)

    return replacer


# ---------------------------------------------------------------------------
# Office / テキスト ファイル処理
# ---------------------------------------------------------------------------

def _process_text_frame(tf, replacer) -> int:
    count = 0
    for para in tf.paragraphs:
        for run in para.runs:
            new = replacer(run.text)
            if new != run.text:
                run.text = new
                count += 1
    return count


def replace_in_pptx(filepath: str, replacer, output_path: str = None) -> int:
    from pptx import Presentation
    prs   = Presentation(filepath)
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                total += _process_text_frame(shape.text_frame, replacer)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        total += _process_text_frame(cell.text_frame, replacer)
    if total:
        prs.save(output_path or filepath)
    elif output_path:
        shutil.copy2(filepath, output_path)
    return total


def replace_in_xlsx(filepath: str, replacer, output_path: str = None) -> int:
    import openpyxl
    wb    = openpyxl.load_workbook(filepath)
    total = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    new = replacer(cell.value)
                    if new != cell.value:
                        cell.value = new
                        total += 1
    if total:
        wb.save(output_path or filepath)
    elif output_path:
        shutil.copy2(filepath, output_path)
    return total


def replace_in_docx(filepath: str, replacer, output_path: str = None) -> int:
    from docx import Document

    def proc_paras(paragraphs):
        n = 0
        for para in paragraphs:
            for run in para.runs:
                new = replacer(run.text)
                if new != run.text:
                    run.text = new
                    n += 1
        return n

    doc   = Document(filepath)
    total = proc_paras(doc.paragraphs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                total += proc_paras(cell.paragraphs)
    if total:
        doc.save(output_path or filepath)
    elif output_path:
        shutil.copy2(filepath, output_path)
    return total


_TXT_ENCODINGS = ("utf-8-sig", "utf-8", "cp932", "latin-1")

def replace_in_txt(filepath: str, replacer, output_path: str = None) -> int:
    text, enc = None, "utf-8"
    for e in _TXT_ENCODINGS:
        try:
            with open(filepath, "r", encoding=e) as f:
                text = f.read()
            enc = e
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if text is None:
        raise ValueError("エンコーディングを判別できませんでした")

    new_text = replacer(text)
    if new_text != text:
        with open(output_path or filepath, "w", encoding=enc) as f:
            f.write(new_text)
        return 1
    elif output_path:
        shutil.copy2(filepath, output_path)
    return 0


HANDLERS = {
    ".pptx": replace_in_pptx,
    ".xlsx": replace_in_xlsx,
    ".docx": replace_in_docx,
    ".txt":  replace_in_txt,
    ".md":   replace_in_txt,   # txt と同じテキスト処理
}


# ---------------------------------------------------------------------------
# pptx → PDF 変換 / スライドマスタ置換（PowerPoint COM）
# ---------------------------------------------------------------------------

def apply_slide_master(pptx_path: str, template_path: str, output_path: str = None) -> str | None:
    """
    PowerPoint COM を使ってスライドマスタをテンプレート (.pptx/.potx) から適用する。
    成功時は None、失敗時はエラーメッセージ文字列を返す。
    """
    try:
        import comtypes.client
    except ImportError:
        return "comtypes が未インストールです (pip install comtypes)"
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            ppt.Visible = 1
            deck = ppt.Presentations.Open(
                os.path.abspath(pptx_path), ReadOnly=False, WithWindow=False)
            deck.ApplyTemplate(os.path.abspath(template_path))
            save_path = os.path.abspath(output_path or pptx_path)
            deck.SaveAs(save_path, 24)   # 24 = ppSaveAsOpenXMLPresentation
            deck.Close()
        finally:
            ppt.Quit()
        return None
    except Exception as e:
        return str(e)


def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> str | None:
    """
    Microsoft PowerPoint COM を使って pptx を PDF に変換する。
    成功時は None、失敗時はエラーメッセージ文字列を返す。
    """
    try:
        import comtypes.client
    except ImportError:
        return "comtypes が未インストールです (pip install comtypes)"
    try:
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            ppt.Visible = 1
            deck = ppt.Presentations.Open(
                os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
            deck.SaveAs(os.path.abspath(pdf_path), 32)   # 32 = ppSaveAsPDF
            deck.Close()
        finally:
            ppt.Quit()
        return None
    except Exception as e:
        return str(e)


# ---------------------------------------------------------------------------
# JSON コメント除去（// スタイル）
# ---------------------------------------------------------------------------

def _strip_json_comments(text: str) -> str:
    result, in_str, i, n = [], False, 0, len(text)
    while i < n:
        c = text[i]
        if c == '\\' and in_str:
            result.append(c); i += 1
            if i < n:
                result.append(text[i]); i += 1
            continue
        if c == '"':
            in_str = not in_str
            result.append(c); i += 1
            continue
        if not in_str and c == '/' and i + 1 < n and text[i + 1] == '/':
            while i < n and text[i] != '\n':
                i += 1
            continue
        result.append(c); i += 1
    return ''.join(result)


# ---------------------------------------------------------------------------
# テキストエリアのパース
# ---------------------------------------------------------------------------

def parse_pairs(text: str) -> tuple[list[tuple[str, str]], list[str]]:
    pairs, errors = [], []
    for i, line in enumerate(text.splitlines(), 1):
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if "=" not in line:
            errors.append(f"行{i}: '=' が見つかりません → {line!r}")
            continue
        old, new = line.split("=", 1)
        if not old:
            errors.append(f"行{i}: 置換前文字が空です")
            continue
        pairs.append((old, new))
    return pairs, errors


def pairs_to_text(pairs: list[dict]) -> str:
    lines = []
    for p in pairs:
        old, new, enabled = p.get("old",""), p.get("new",""), p.get("enabled", True)
        line = f"{old}={new}"
        if not enabled:
            line = f"# {line}"
        lines.append(line)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# メインアプリ
# ---------------------------------------------------------------------------

class ReplaceAllApp:
    def __init__(self, root: tk.Tk):
        self.root = root
        self.root.title("一括文字置換ツール")
        self.root.geometry("1280x720")
        self.root.resizable(True, True)
        self._build_ui()

    # ---- UI 構築 -----------------------------------------------------------

    def _build_ui(self):
        p = {"padx": 6, "pady": 3}

        # ── 2列メインフレーム ──────────────────────────────────────────────────
        main = ttk.Frame(self.root)
        main.pack(fill=tk.BOTH, expand=True, padx=6, pady=4)

        # 左列：オプション群（固定幅）
        left = ttk.Frame(main, width=550)
        left.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 6))
        left.pack_propagate(False)

        # 右列：置換ペア・ログ（可変幅）
        right = ttk.Frame(main)
        right.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # ══════════════════════════════════════════
        # 左列
        # ══════════════════════════════════════════

        # === フォルダ選択 ===
        frm_folder = ttk.LabelFrame(left, text="対象フォルダ")
        frm_folder.pack(fill=tk.X, **p)

        self.folder_var = tk.StringVar()
        ttk.Entry(frm_folder, textvariable=self.folder_var).pack(
            side=tk.LEFT, padx=4, pady=4, fill=tk.X, expand=True)
        ttk.Button(frm_folder, text="参照...", command=self._browse_folder).pack(
            side=tk.LEFT, padx=4, pady=4)

        # === 対象オプション ===
        frm_opt = ttk.Frame(left)
        frm_opt.pack(fill=tk.X, **p)

        self.recurse_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(frm_opt, text="サブフォルダを含める",
                        variable=self.recurse_var).pack(side=tk.LEFT)

        self.ext_vars: dict[str, tk.BooleanVar] = {}
        for _ext in (".pptx", ".xlsx", ".docx", ".txt", ".md"):
            v = tk.BooleanVar(value=True)
            self.ext_vars[_ext] = v
            ttk.Checkbutton(frm_opt, text=_ext, variable=v).pack(side=tk.LEFT, padx=4)

        # === マッチングオプション ===
        frm_match = ttk.LabelFrame(left, text="マッチングオプション")
        frm_match.pack(fill=tk.X, **p)

        self.case_var     = tk.BooleanVar(value=False)
        self.fullhalf_var = tk.BooleanVar(value=False)

        ttk.Checkbutton(frm_match, text="大文字・小文字を区別する",
                        variable=self.case_var).pack(side=tk.LEFT, padx=8, pady=4)
        ttk.Checkbutton(frm_match, text="全角・半角を区別する",
                        variable=self.fullhalf_var).pack(side=tk.LEFT, padx=8, pady=4)

        # === 出力オプション ===
        frm_out = ttk.LabelFrame(left, text="出力オプション")
        frm_out.pack(fill=tk.X, **p)

        frm_out1 = ttk.Frame(frm_out)
        frm_out1.pack(fill=tk.X, padx=4, pady=(4, 2))

        self.out_separate_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(frm_out1, text="別フォルダに出力する（元ファイルを変更しない）",
                        variable=self.out_separate_var,
                        command=self._on_out_toggle).pack(side=tk.LEFT)

        frm_out2 = ttk.Frame(frm_out)
        frm_out2.pack(fill=tk.X, padx=4, pady=(0, 2))

        ttk.Label(frm_out2, text="出力先:").pack(side=tk.LEFT)
        self.out_dir_var = tk.StringVar()
        self._out_entry = ttk.Entry(frm_out2, textvariable=self.out_dir_var)
        self._out_entry.pack(side=tk.LEFT, padx=(4, 2), fill=tk.X, expand=True)
        self._out_browse_btn = ttk.Button(frm_out2, text="参照...",
                                          command=self._browse_out_folder, width=7)
        self._out_browse_btn.pack(side=tk.LEFT, padx=2)

        frm_out3 = ttk.Frame(frm_out)
        frm_out3.pack(fill=tk.X, padx=22, pady=(0, 6))

        self.out_structure_var = tk.BooleanVar(value=False)
        self._out_struct_cb = ttk.Checkbutton(
            frm_out3, text="フォルダ構成を維持する（未チェック＝フラット出力）",
            variable=self.out_structure_var)
        self._out_struct_cb.pack(side=tk.LEFT)

        self._on_out_toggle()

        # === ファイル名変更オプション ===
        frm_fname = ttk.LabelFrame(left, text="ファイル名変更オプション")
        frm_fname.pack(fill=tk.X, **p)

        frm_fname1 = ttk.Frame(frm_fname)
        frm_fname1.pack(fill=tk.X, padx=4, pady=(4, 2))

        self.fname_replace_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(frm_fname1, text="置換ペアをファイル名にも適用する",
                        variable=self.fname_replace_var).pack(side=tk.LEFT)

        frm_fname2 = ttk.Frame(frm_fname)
        frm_fname2.pack(fill=tk.X, padx=4, pady=(2, 6))

        self.fname_add_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(frm_fname2, text="ファイル名に文字列を追加する",
                        variable=self.fname_add_var,
                        command=self._on_fname_add_toggle).pack(side=tk.LEFT)

        self.fname_pos_var = tk.StringVar(value="prefix")
        self._rb_prefix = ttk.Radiobutton(frm_fname2, text="冒頭",
                                          variable=self.fname_pos_var, value="prefix")
        self._rb_suffix = ttk.Radiobutton(frm_fname2, text="末尾",
                                          variable=self.fname_pos_var, value="suffix")
        self._rb_prefix.pack(side=tk.LEFT, padx=(10, 2))
        self._rb_suffix.pack(side=tk.LEFT, padx=(0, 4))

        self.fname_str_var = tk.StringVar(value="【マスキング済】")
        self._fname_entry = ttk.Entry(frm_fname2, textvariable=self.fname_str_var, width=22)
        self._fname_entry.pack(side=tk.LEFT, padx=2)

        self._on_fname_add_toggle()

        # === 変換オプション ===
        frm_conv = ttk.LabelFrame(left, text="変換オプション")
        frm_conv.pack(fill=tk.X, **p)

        frm_conv1 = ttk.Frame(frm_conv)
        frm_conv1.pack(fill=tk.X, padx=4, pady=(4, 2))

        self.pdf_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            frm_conv1,
            text="pptx を PDF に変換する  ※ PowerPoint 必要",
            variable=self.pdf_var,
        ).pack(side=tk.LEFT, padx=4)

        frm_conv2 = ttk.Frame(frm_conv)
        frm_conv2.pack(fill=tk.X, padx=4, pady=(2, 2))

        self.master_var = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            frm_conv2,
            text="スライドマスタを置換する  ※ PowerPoint 必要",
            variable=self.master_var,
            command=self._on_master_toggle,
        ).pack(side=tk.LEFT, padx=4)

        frm_conv3 = ttk.Frame(frm_conv)
        frm_conv3.pack(fill=tk.X, padx=22, pady=(0, 6))

        ttk.Label(frm_conv3, text="テンプレート:").pack(side=tk.LEFT)
        self.master_tmpl_var = tk.StringVar()
        self._master_entry = ttk.Entry(frm_conv3, textvariable=self.master_tmpl_var)
        self._master_entry.pack(side=tk.LEFT, padx=(4, 2), fill=tk.X, expand=True)
        self._master_browse_btn = ttk.Button(frm_conv3, text="参照...",
                                             command=self._browse_template, width=7)
        self._master_browse_btn.pack(side=tk.LEFT, padx=2)

        self._on_master_toggle()

        # ══════════════════════════════════════════
        # 右列
        # ══════════════════════════════════════════

        # === 置換ペア ===
        frm_pairs = ttk.LabelFrame(
            right, text="置換ペア  （書式: 置換前=置換後、1行1組、# でコメント）")
        frm_pairs.pack(fill=tk.BOTH, expand=True, **p)

        self.pairs_text = tk.Text(frm_pairs, font=("Consolas", 10), undo=True)
        sb = ttk.Scrollbar(frm_pairs, command=self.pairs_text.yview)
        self.pairs_text.config(yscrollcommand=sb.set)
        sb.pack(side=tk.RIGHT, fill=tk.Y)
        self.pairs_text.pack(fill=tk.BOTH, expand=True, padx=4, pady=4)
        self.pairs_text.insert("1.0", "# 例: 旧文字=新文字\n")

        frm_preset = ttk.Frame(frm_pairs)
        frm_preset.pack(fill=tk.X, padx=4, pady=(0, 4))

        ttk.Button(frm_preset, text="プリセット保存",
                   command=self._save_preset, width=14).pack(side=tk.LEFT, padx=2)
        ttk.Button(frm_preset, text="プリセット読込",
                   command=self._load_preset, width=14).pack(side=tk.LEFT, padx=2)
        ttk.Button(frm_preset, text="全クリア",
                   command=self._clear_pairs, width=9).pack(side=tk.LEFT, padx=2)

        # === 実行ボタン群 ===
        frm_run = ttk.Frame(right)
        frm_run.pack(fill=tk.X, **p)

        self.run_btn = ttk.Button(frm_run, text="▶ 置換実行",
                                  command=self._run_replacement)
        self.run_btn.pack(side=tk.LEFT, padx=4)
        ttk.Button(frm_run, text="ログをクリア",
                   command=self._clear_log).pack(side=tk.LEFT, padx=4)
        ttk.Button(frm_run, text="ログを保存...",
                   command=self._save_log).pack(side=tk.LEFT, padx=4)

        self.status_var = tk.StringVar(value="待機中")
        ttk.Label(frm_run, textvariable=self.status_var,
                  foreground="gray").pack(side=tk.RIGHT, padx=8)

        # === ログ ===
        frm_log = ttk.LabelFrame(right, text="処理ログ")
        frm_log.pack(fill=tk.BOTH, expand=True, **p)

        self.log_text = scrolledtext.ScrolledText(
            frm_log, font=("Consolas", 9), state=tk.DISABLED, height=8)
        self.log_text.pack(fill=tk.BOTH, expand=True, padx=4, pady=4)

        self.log_text.tag_config("ok",    foreground="#006600")
        self.log_text.tag_config("skip",  foreground="#888800")
        self.log_text.tag_config("error", foreground="#cc0000")
        self.log_text.tag_config("info",  foreground="#000099")
        self.log_text.tag_config("head",  foreground="#000000",
                                  font=("Consolas", 9, "bold"))

    # ---- 出力オプション トグル -----------------------------------------------

    def _on_out_toggle(self):
        enabled = self.out_separate_var.get()
        state   = tk.NORMAL if enabled else tk.DISABLED
        self._out_entry.config(state=state)
        self._out_browse_btn.config(state=state)
        self._out_struct_cb.config(state=state)

    def _browse_out_folder(self):
        folder = filedialog.askdirectory(title="出力先フォルダを選択")
        if folder:
            self.out_dir_var.set(folder)

    # ---- ファイル名オプション トグル ------------------------------------------

    def _on_fname_add_toggle(self):
        state = tk.NORMAL if self.fname_add_var.get() else tk.DISABLED
        self._rb_prefix.config(state=state)
        self._rb_suffix.config(state=state)
        self._fname_entry.config(state=state)

    # ---- 変換オプション トグル -----------------------------------------------

    def _on_master_toggle(self):
        state = tk.NORMAL if self.master_var.get() else tk.DISABLED
        self._master_entry.config(state=state)
        self._master_browse_btn.config(state=state)

    def _browse_template(self):
        path = filedialog.askopenfilename(
            title="テンプレートファイルを選択（.pptx / .potx）",
            filetypes=[
                ("PowerPoint ファイル", "*.pptx *.potx"),
                ("すべてのファイル", "*.*"),
            ],
        )
        if path:
            self.master_tmpl_var.set(path)

    # ---- ペア操作 ----------------------------------------------------------

    def _clear_pairs(self):
        self.pairs_text.delete("1.0", tk.END)

    # ---- プリセット --------------------------------------------------------

    def _save_preset(self):
        raw = self.pairs_text.get("1.0", tk.END)
        pairs_raw, errors = parse_pairs(raw)
        if errors:
            if not messagebox.askyesno("入力エラー",
                                       "\n".join(errors) + "\n\nエラー行を除いて保存しますか？"):
                return

        path = filedialog.asksaveasfilename(
            title="プリセットを保存",
            defaultextension=".json",
            filetypes=[("JSON ファイル", "*.json"), ("すべてのファイル", "*.*")],
            initialfile="preset.json",
        )
        if not path:
            return

        all_pairs = []
        for line in raw.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                inner = stripped[1:].strip()
                if "=" in inner:
                    o, n = inner.split("=", 1)
                    if o:
                        all_pairs.append({"enabled": False, "old": o, "new": n})
            elif "=" in stripped:
                o, n = stripped.split("=", 1)
                if o:
                    all_pairs.append({"enabled": True, "old": o, "new": n})

        data = {
            "version": 1,
            "case_sensitive":     self.case_var.get(),
            "fullhalf_sensitive": self.fullhalf_var.get(),
            "pairs": all_pairs,
        }
        try:
            with open(path, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            messagebox.showinfo("保存完了", f"プリセットを保存しました:\n{path}")
        except Exception as e:
            messagebox.showerror("保存エラー", str(e))

    def _load_preset(self):
        path = filedialog.askopenfilename(
            title="プリセットを読込",
            filetypes=[("JSON ファイル", "*.json"), ("すべてのファイル", "*.*")],
        )
        if not path:
            return
        try:
            with open(path, "r", encoding="utf-8") as f:
                raw_text = f.read()
            data = json.loads(_strip_json_comments(raw_text))
        except json.JSONDecodeError as e:
            messagebox.showerror("JSON 解析エラー", f"JSON の形式が正しくありません:\n{e}")
            return
        except Exception as e:
            messagebox.showerror("読込エラー", str(e))
            return

        self.case_var.set(data.get("case_sensitive", False))
        self.fullhalf_var.set(data.get("fullhalf_sensitive", False))

        text = pairs_to_text(data.get("pairs", []))
        self.pairs_text.delete("1.0", tk.END)
        self.pairs_text.insert("1.0", text)

    # ---- ログ --------------------------------------------------------------

    def _log(self, message: str, tag: str = ""):
        ts = datetime.now().strftime("%H:%M:%S")
        self.log_text.config(state=tk.NORMAL)
        self.log_text.insert(tk.END, f"[{ts}] {message}\n", tag)
        self.log_text.see(tk.END)
        self.log_text.config(state=tk.DISABLED)

    def _clear_log(self):
        self.log_text.config(state=tk.NORMAL)
        self.log_text.delete("1.0", tk.END)
        self.log_text.config(state=tk.DISABLED)

    def _save_log(self):
        path = filedialog.asksaveasfilename(
            title="ログを保存",
            defaultextension=".txt",
            filetypes=[("テキストファイル", "*.txt"), ("すべてのファイル", "*.*")],
            initialfile=f"replace_log_{datetime.now():%Y%m%d_%H%M%S}.txt",
        )
        if not path:
            return
        try:
            with open(path, "w", encoding="utf-8") as f:
                f.write(self.log_text.get("1.0", tk.END))
            messagebox.showinfo("保存完了", f"ログを保存しました:\n{path}")
        except Exception as e:
            messagebox.showerror("保存エラー", str(e))

    # ---- フォルダ参照 ------------------------------------------------------

    def _browse_folder(self):
        folder = filedialog.askdirectory(title="対象フォルダを選択")
        if folder:
            self.folder_var.set(folder)

    # ---- 置換実行 ----------------------------------------------------------

    def _run_replacement(self):
        folder = self.folder_var.get().strip()
        if not folder or not os.path.isdir(folder):
            messagebox.showerror("エラー", "有効なフォルダを選択してください")
            return

        raw = self.pairs_text.get("1.0", tk.END)
        pairs_raw, errors = parse_pairs(raw)
        if errors:
            messagebox.showerror("入力エラー", "\n".join(errors))
            return

        case_sensitive     = self.case_var.get()
        fullhalf_sensitive = self.fullhalf_var.get()
        # 置換ペア未入力の場合は恒等関数（コピー・PDF変換のみ実行）
        replacer = build_combined_replacer(pairs_raw, case_sensitive, fullhalf_sensitive) \
                   or (lambda text: text)

        exts = [e for e, v in self.ext_vars.items() if v.get()]
        if not exts:
            messagebox.showerror("エラー", "対象拡張子を1つ以上選択してください")
            return

        # 別フォルダ出力の検証
        out_separate = self.out_separate_var.get()
        out_dir      = self.out_dir_var.get().strip()
        if out_separate and not out_dir:
            messagebox.showerror("エラー", "出力先フォルダを指定してください")
            return

        # スライドマスタ置換の検証
        master_apply    = self.master_var.get()
        master_template = self.master_tmpl_var.get().strip()
        if master_apply and not master_template:
            messagebox.showerror("エラー", "スライドマスタ置換のテンプレートファイルを指定してください")
            return
        if master_apply and not os.path.isfile(master_template):
            messagebox.showerror("エラー", f"テンプレートファイルが見つかりません:\n{master_template}")
            return

        self.run_btn.config(state=tk.DISABLED)
        self.status_var.set("処理中...")
        self.root.update_idletasks()

        opts = {
            "case_sensitive":     case_sensitive,
            "fullhalf_sensitive": fullhalf_sensitive,
            "recurse":            self.recurse_var.get(),
            # 出力
            "out_separate":       out_separate,
            "out_dir":            out_dir,
            "out_structure":      self.out_structure_var.get(),
            # ファイル名
            "fname_replace":      self.fname_replace_var.get(),
            "fname_add":          self.fname_add_var.get(),
            "fname_position":     self.fname_pos_var.get(),
            "fname_string":       self.fname_str_var.get(),
            # 変換
            "pdf_convert":        self.pdf_var.get(),
            "master_apply":       master_apply,
            "master_template":    master_template,
        }

        t = threading.Thread(
            target=self._do_replacement,
            args=(folder, replacer, pairs_raw, exts, opts),
            daemon=True,
        )
        t.start()

    def _do_replacement(self, folder, replacer, pairs_raw, exts, opts):
        start = datetime.now()

        out_separate    = opts["out_separate"]
        out_dir         = opts["out_dir"]
        out_structure   = opts["out_structure"]
        fname_replace   = opts["fname_replace"]
        fname_add       = opts["fname_add"]
        fname_position  = opts["fname_position"]
        fname_string    = opts["fname_string"]
        pdf_convert     = opts["pdf_convert"]
        master_apply    = opts["master_apply"]
        master_template = opts["master_template"]
        recurse         = opts["recurse"]

        opt_case = "区別する" if opts["case_sensitive"] else "区別しない"
        opt_fh   = "区別する" if opts["fullhalf_sensitive"] else "区別しない"

        self.root.after(0, lambda: self._log(f"=== 置換開始: {folder} ===", "head"))
        self.root.after(0, lambda: self._log(
            f"大文字小文字: {opt_case} / 全角半角: {opt_fh}", "info"))
        self.root.after(0, lambda: self._log(
            f"置換ペア ({len(pairs_raw)}組): " +
            " / ".join(f'"{o}"→"{n}"' for o, n in pairs_raw), "info"))
        if out_separate:
            mode = "構成維持" if out_structure else "フラット"
            self.root.after(0, lambda: self._log(
                f"出力先: {out_dir}  [{mode}]", "info"))
        if pdf_convert:
            self.root.after(0, lambda: self._log("PDF変換: 有効", "info"))

        total_files   = 0
        changed_files = 0
        total_reps    = 0
        renamed_files = 0
        pdf_count     = 0
        master_count  = 0
        error_count   = 0

        walker = os.walk(folder) if recurse else [(folder, [], os.listdir(folder))]

        for dirpath, _, filenames in walker:
            for fname in sorted(filenames):
                stem, ext = os.path.splitext(fname)
                if ext.lower() not in exts:
                    continue
                total_files += 1
                fpath = os.path.join(dirpath, fname)
                rel   = os.path.relpath(fpath, folder)

                # --- 出力パスの決定 ---
                if out_separate and out_dir:
                    if out_structure:
                        out_path = os.path.join(out_dir, rel)
                    else:
                        out_path = os.path.join(out_dir, fname)
                    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
                else:
                    out_path = None   # 上書き

                # --- 内容の置換 ---
                handler = HANDLERS[ext.lower()]
                try:
                    n = handler(fpath, replacer, out_path)
                    if n:
                        changed_files += 1
                        total_reps    += n
                        msg, tag = f"[変更]    {rel}  ({n}箇所)", "ok"
                    else:
                        msg, tag = f"[変更なし] {rel}", "skip"
                except Exception as e:
                    error_count += 1
                    msg, tag = f"[エラー]  {rel}: {e}", "error"

                self.root.after(0, lambda m=msg, t=tag: self._log(m, t))

                # 以降のファイル名操作・PDF変換は実際のファイルパスに対して行う
                target_path = out_path or fpath

                # --- スライドマスタ置換（pptx のみ）---
                if master_apply and ext.lower() == ".pptx":
                    err = apply_slide_master(target_path, master_template)
                    if err:
                        error_count += 1
                        mmsg = f"[マスタ置換エラー] {rel}: {err}"
                        self.root.after(0, lambda m=mmsg: self._log(m, "error"))
                    else:
                        master_count += 1
                        mmsg = f"[マスタ置換]  {os.path.basename(target_path)}"
                        self.root.after(0, lambda m=mmsg: self._log(m, "ok"))

                # --- ファイル名の計算（PDF変換前に確定させる）---
                new_stem = stem
                if fname_replace:
                    new_stem = replacer(new_stem)
                if fname_add and fname_string:
                    if fname_position == "prefix":
                        if not new_stem.startswith(fname_string):
                            new_stem = fname_string + new_stem
                    else:
                        if not new_stem.endswith(fname_string):
                            new_stem = new_stem + fname_string

                # --- PDF変換（pptx のみ）---
                pdf_converted = False
                if pdf_convert and ext.lower() == ".pptx":
                    pdf_path = os.path.join(os.path.dirname(target_path), new_stem + ".pdf")
                    err = convert_pptx_to_pdf(target_path, pdf_path)
                    if err:
                        error_count += 1
                        pmsg = f"[PDF変換エラー] {rel}: {err}"
                        self.root.after(0, lambda m=pmsg: self._log(m, "error"))
                    else:
                        pdf_count += 1
                        pdf_converted = True
                        pmsg = f"[PDF変換]  {os.path.relpath(pdf_path, folder if not out_separate else out_dir)}"
                        self.root.after(0, lambda m=pmsg: self._log(m, "ok"))
                        try:
                            os.remove(target_path)
                        except Exception as e:
                            error_count += 1
                            rmsg = f"[pptx削除エラー] {os.path.basename(target_path)}: {e}"
                            self.root.after(0, lambda m=rmsg: self._log(m, "error"))

                # --- ファイル名の変更（PDF変換済みの場合はpptxが削除済みのためスキップ）---
                if not pdf_converted and new_stem != stem:
                    new_fname = new_stem + ext
                    new_target = os.path.join(os.path.dirname(target_path), new_fname)
                    try:
                        if os.path.exists(new_target):
                            rmsg = f"[名前変更スキップ] {rel} → {new_fname}（同名ファイルが存在）"
                            self.root.after(0, lambda m=rmsg: self._log(m, "error"))
                        else:
                            os.rename(target_path, new_target)
                            renamed_files += 1
                            rmsg = f"[名前変更]  {os.path.basename(target_path)} → {new_fname}"
                            self.root.after(0, lambda m=rmsg: self._log(m, "ok"))
                    except Exception as e:
                        error_count += 1
                        rmsg = f"[名前変更エラー] {rel}: {e}"
                        self.root.after(0, lambda m=rmsg: self._log(m, "error"))

        elapsed = (datetime.now() - start).total_seconds()
        summary = (
            f"=== 完了 ({elapsed:.1f}秒)  "
            f"対象:{total_files}件  内容変更:{changed_files}件  "
            f"置換:{total_reps}箇所  名前変更:{renamed_files}件  "
            f"PDF変換:{pdf_count}件  マスタ置換:{master_count}件  エラー:{error_count}件 ==="
        )
        self.root.after(0, lambda: self._log(summary, "head"))
        self.root.after(0, lambda: self.status_var.set(
            f"完了 — 変更:{changed_files}/{total_files}  名前:{renamed_files}  PDF:{pdf_count}  マスタ:{master_count}"))
        self.root.after(0, lambda: self.run_btn.config(state=tk.NORMAL))


# ---------------------------------------------------------------------------
# エントリポイント
# ---------------------------------------------------------------------------

def _check_deps():
    missing = []
    for pkg, mod in [("python-pptx", "pptx"),
                     ("openpyxl",    "openpyxl"),
                     ("python-docx", "docx")]:
        try:
            __import__(mod)
        except ImportError:
            missing.append(pkg)
    return missing


if __name__ == "__main__":
    missing = _check_deps()
    if missing:
        print("以下のライブラリが不足しています。インストールしてください:")
        print(f"  pip install {' '.join(missing)}")
        sys.exit(1)

    root = tk.Tk()
    ReplaceAllApp(root)
    root.mainloop()
